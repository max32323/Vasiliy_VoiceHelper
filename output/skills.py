"""ИМПОРТ БИБЛИОТЕК"""
import sqlite3
import string
import webbrowser, sys, subprocess, voice, random, pyperclip
import nltk
import openai
import pygame
import win32api
import win32gui
import password_menu
import datetime
import os
from pathlib import Path
from cryptography.fernet import Fernet
import winshell
import psutil
from gnewsclient import gnewsclient
from pptx import Presentation
import requests
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume

connect_path = sqlite3.connect('user_path.db')
cursor_path = connect_path.cursor()
cursor_path.execute("""SELECT * FROM your_path""")
path = cursor_path.fetchall()

"""ФУНКЦИИ ДЛЯ УПРОЩЕНИЯ"""
b = True

def russian():
    russian_letters = ''.join(chr(i) for i in range(1040, 1104))  # А-Я, а-я
    digits = string.digits
    special_characters = '1234567890!@#$%^&*()-—_=+№%[]{},.<>?;:"\\|\\\\`~–' + "'" + ' '
    result_string = russian_letters + digits + special_characters
    return result_string
def time_all():
    time = datetime.datetime.now()
    hour = time.hour
    minute = time.minute
    day = time.day
    month_dist = {1: 'Январь', 2: 'Февраль', 3: 'Март', 4: 'Апрель', 5: 'Май', 6: 'Июнь', 7: 'Июль',
                  8: 'Август', 9: 'Сентябрь', 10: 'Октябрь', 11: 'Ноябрь', 12: 'Декабрь'}
    for key, values in month_dist.items():
        if time.month == key:
            month = values
    if int(hour) < 10:
        hour = '0' + str(hour)
    if int(minute) < 10:
        minute = '0' + str(minute)
    if int(day) < 10:
        day = '0' + str(day)
    return f'{time.year}: {day} {month} - {hour}:{minute}'

"""ОСНОВНЫЕ ФУНКЦИИ"""

def setCyrillicLayout():
    window_handle = win32gui.GetForegroundWindow()
    result = win32api.SendMessage(window_handle, 0x0050, 0, 0x04190419)

def setEngLayout():
    window_handle = win32gui.GetForegroundWindow()
    result = win32api.SendMessage(window_handle, 0x0050, 0, 0x04090409)

def browser():
    webbrowser.open('https://ya.ru')
    voice.speaker('Открываю браузер')

def browser_youtube():
    webbrowser.open('https://www.youtube.com')
    voice.speaker('открываю ютуб')

def open_telegramm():
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][1]}')
        voice.speaker(random.choice(['открываю для вас телеграмм', 'запускаю телегу', 'запускаю, удачных переписок!', 'открываю мессенджер']))
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def game_steam():
    try:
        voice.speaker(random.choice(['открываю стим', 'запускаю каталог игр', 'запустил, удачных игр!']))
        subprocess.Popen('C:\Program Files (x86)\Steam\steam.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def jokes():
    jokes = [
        "как называется трава в клубе? клубника.",
        "что говорит кекс в кофейне? не кекси!",
        "как зовут короля, который никогда не дуется? польский.",
        "почему мороженое всегда счастливо? потому что оно всегда мокрое.",
        "что должен сделать дом, когда ему становится невыносимо? снести крышу.",
        "как называется самый непослушный овощ? помидорка.",
        "почему рыба не спорит? потому что она рыба.",
        'как называют человека, который продал свою печень? обеспеченный.',
        "Как называется процесс перевода кошачьего языка? - Переводение!",
        "Почему программисты не любят играть в футбол? - Потому что ворота всегда закрыты.",
        "Что говорит телефон программисту? - Помогите, меня программисты убивают!",
        "Как программист купил стул? - Через интернет.",
        "Почему алгоритмы не пьют кофе? - Потому что им не нужен java.",
        "Что сказал массив строке? - Привет, строка!",
        "Зачем программисты так любят гречку? Потому что в ней много библиотек!",
        "Как программист пересекает реку? По памяти!",
        "Почему программисты не ходят в школу? Потому что им уже все if и else известно!",
        "Что говорит начинающий программист, когда видит баг? Hello, world!",
        "Как программист перекрывает огород? Циклом for!",
        "Почему программисты так медленно реагируют на шутки? Они сначала обрабатывают!",
        "Что делает кот, когда программист пишет код? CTRL+C, CTRL+V!",
        "Какой метод программирования самый популярный? Подход сверху вниз, белка вверх ногами!",
        "Почему программисты не могут войти в библиотеку? Потому что пароль неверный!",
        "Чем программисты празднуют новый год? Начиная с 01.01!"
    ]
    random_joke = random.choice(jokes)
    voice.speaker(random_joke)


def data_new():
    time_new = time_all()
    voice.speaker(f'Дата на сегодня: {time_new}')


def menu_weather(city: str):
    voice.speaker(random.choice(['узнаю погоду подождите', 'ищу прогноз погоды на сегодня', 'начал поиск погоды', 'начал поиск погоды', 'сейчас расскажу']))
    current_datetime = datetime.datetime.now()
    moun = str(current_datetime).split()[0]
    conn = sqlite3.connect('wether_data.db')

    cur = conn.cursor()

    cur.execute("""
        CREATE TABLE IF NOT EXISTS weather (
        city TEXT,
        data TEXT,
        time TEXT,
        temperature REAL,
        weather TEXT
        )
    """)
    conn.commit()
    cur.execute("""SELECT * FROM weather""")
    all_weather = cur.fetchall()
    print(len(all_weather))
    if len(all_weather) >= 10:
        cur.execute("""DELETE FROM weather""")
        conn.commit()
    time = datetime.datetime.now()
    time_new = str(time.hour)+':'+str(time.minute)

    text_new = city.split()
    text_new_2 = text_new[5:]
    city_new = ''
    if len(text_new_2) == 1:
        for i in text_new_2:
            city_new = city_new + i
    else:
        for i in text_new_2:
            city_new += i
            city_new += ' '
    print(city_new)
    try:
        url = 'https://api.openweathermap.org/data/2.5/weather?q=' + city_new + '&units=metric&lang=ru&appid=79d1ca96933b0328e1c7e3e7a26cb347'
        weather_data = requests.get(url).json()
        temp = round(weather_data['main']['temp'])
        weather_status = weather_data['weather'][0]['description']
        if int(temp) <= 10:
            voice.speaker(
                f'Сейчас в городе {city_new} температура состовляет {round(temp)} градус цельсия. Лучше оденьтесь потеплее. Тип погоды: {weather_status}')
        if int(temp) <= 20 and int(temp) >= 11:
            voice.speaker(
                f'Сейчас в городе {city_new} температура состовляет {round(temp)} градус цельсия. Неплохая погода для прогулки. Тип погоды: {weather_status}')
        elif int(temp) > 20:
            voice.speaker(
                f'Сейчас в городе {city_new} температура состовляет {round(temp)} градус цельсия. Сегодня будет очень жарко. Тип погоды: {weather_status}')
        cur.execute("""INSERT INTO weather VALUES (?,?,?,?,?)""",
                    (str(city_new), str(moun), str(time_new), temp, str(weather_status)))
        conn.commit()
    except:
        voice.speaker('Извините я не смог найти погоду. Скажите мне существующий город или скажите мне в именительном падеже. Также проверте соединение с интернетом')





def offpc():
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
    t = datetime.datetime.now()
    h = t.hour
    if h >= 21:
        voice.speaker(f'Выключаю ваш компьютер! Спокойной ночи {name}!')
    else:
        voice.speaker(f'Выключаю ваш компьютер! До свидания {name}')
    if os.name == 'posix':
        os.system('shutdown -h now')
    elif os.name == 'nt':
        os.system('shutdown /s /t 0')
    else:
        voice.speaker("Данная операционная система не поддерживается")

def offbot():
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    l = ['до свидание', 'до скорых встреч', 'ещё увидимся', 'пока']
    l_random = random.choice(l)
    time_all_1 = datetime.datetime.now()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        if (time_all_1.hour >= 21 and time_all_1.hour <= 23) or (time_all_1.hour >= 0 and time_all_1.hour <= 4):
            l2 = ['добрый снов', 'спокойной ночи', 'до завтра']
            l_random2 = random.choice(l2)
            if random.randint(1,2) == 1:
                voice.speaker(f'{l_random2}, {name}!')
            else:
                voice.speaker(f'Спасибо за работу! {l_random2}, {name}!')
        else:
            if random.randint(1,2) == 1:
                voice.speaker(f'{l_random}, {name}!')
            else:
                voice.speaker(f'Спасибо за работу! {l_random}, {name}!')
    sys.exit()


# def menu_password():
#     try:
#         cursor_path.execute("""SELECT * FROM your_path""")
#         path = cursor_path.fetchall()
#         subprocess.Popen(f'{path[0][0]}\output\password_menu.exe')
#     except:
#         voice.speaker('Ошибка! Проверьте расположения файла')

def menu_password_v2():
    password_menu.generate_v2()
    voice.speaker('готово! я сгенерировал и сохранил пароль!')

def passive():
    pass
def time_now():
    n = datetime.datetime.now()
    voice.speaker(f'Сейчас на улице {n.hour} часов {n.minute} минут по МСК')

def learning_programmer():
    webbrowser.open('https://start.1t.ru/user')
    voice.speaker('открываю браузер')

def version():
    VOICEHELPER_VERSION = '2.1.0'
    voice.speaker(f'Моя версия: {VOICEHELPER_VERSION}. Самая последняя версия на данный момент.')

def lern():
    voice.speaker('открываю переводчик')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\menu_translate.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

# def voice_name():
#     voice.engine.setProperty('voice', voice.voices[2].id)
#
# def voice_return():
#     voice.engine.setProperty('voice', voice.voices[3].id)

def name():
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        time_user = datetime.datetime.now()
        if time_user.hour >= 6 and time_user.hour <= 12:
            voice.speaker(f'доброе утро {name}')
        elif time_user.hour >= 12 and time_user.hour <= 17:
            hi = random.choice(['привет', 'здравствуйте', 'приветствую вас', 'алоха'])
            voice.speaker(f'{hi} {name}')
        else:
            voice.speaker(f'Добрый вечер {name}')

def shoping():
    voice.speaker('открываю меню составления списка покупок')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\shop.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def disp():
    voice.speaker('открываю диспечер задач')
    os.system('taskmgr')

def provodnic():
    voice.speaker('открываю проводник')
    subprocess.run(["explorer"])

def vk():
    voice.speaker('открываю вконтакте')
    webbrowser.open('https://vk.com/')

def name_write(name_user):
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    file = open(f'{path[0][0]}\my_name', 'w', encoding='utf-8')
    file.write(name_user)
    file.close()

def play():
    voice.speaker('подбрасываю монетку')
    r = random.randint(1,2)
    if r == 1:
        voice.speaker('выпала решка!')
    elif r == 2:
        voice.speaker('выпал орёл!')

def tren():
    voice.speaker('открываю меню тренировки печати')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\peth.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def mine():
    voice.speaker('Хорошо, приятной вам игры!')
    try:
        executable_path = Path(
            'c:/users/user/appdata/roaming/.minecraft/klauncher.exe')  # можно выбрать другой путь игры
        os.startfile(executable_path.resolve())
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')
def write_zamet(text: str):
    if 'запиши заметки' in text or 'забежав заметки' in text:
        text_new = text.split()
        text_new2 = text_new[3:]
    else:
        text_new = text.split()
        text_new2 = text_new[4:]
    normal_text = ''
    for text_el in text_new2:
        normal_text = normal_text + text_el
        normal_text = normal_text + ' '
    time_now = time_all()
    print(normal_text)
    voice.speaker('сохранено в теиксти файл записки')
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\заметки и шифры\записки.txt', 'a', encoding='utf-8') as file:
        file.write(f'ЗАМЕТКА: {str(normal_text).capitalize()}\nДАТА СОХРАНЕНИЯ: {time_now}\n\n')

def remove_zamet():
    voice.speaker('очищаю заметки')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        f = open(f'{path[0][0]}\заметки и шифры\записки.txt', 'w')
        f.close()
    except:
        voice.speaker('Извините, файл отсутствует')

# def meat():
#     voice.speaker('всё отлично, а у вас')

def shifr(text: str):
    try:
        voice.speaker('готово! текст зашифрован и сохранён в буфет обмена')
        if 'зашифрованный текст' in text:
            text_new = text.split()
            text_new2 = text_new[3:]
        elif 'зажав рой мой текст' in text:
            text_new = text.split()
            text_new2 = text_new[5:]
        else:
            text_new = text.split()
            text_new2 = text_new[4:]
        text_norm1 = ''
        for text_norm in text_new2:
            text_norm1 = text_norm1 + text_norm
            text_norm1 = text_norm1 + ' '
        text_norm1 = text_norm1.strip()
        key = Fernet.generate_key()
        cur = Fernet(key)
        text_new3 = cur.encrypt(text_norm1.encode())
        pyperclip.copy(str(text_new3))
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        with open(f'{path[0][0]}\заметки и шифры\шифры.txt', 'a', encoding='utf-8') as file:
            file.write(f'СООБЩЕНИЕ: {text_norm1}\n')
            file.write(f'ШИФРОВОНОЕ СООБЩЕНИЕ: {text_new3}\n')
    except:
        voice.speaker('Ошибка! Попробуйте проверить наличие файла шифры для шифрации текста')


def clean2():
    voice.speaker('готово. Шифры были удалены')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        f = open(f'{path[0][0]}\заметки и шифры\шифры.txt', 'w')
        f.close()
    except:
        voice.speaker('Ошибка выполнения! Файл отсутствует!')

# def veb():
#     webbrowser.open('https://start.1t.ru/user/webinar')


def open_web(text):
    try:
        voice.speaker('открываю сайт')
        text_new1 = text.split()
        text_new2 = text_new1[3:]
        answer = ''
        for word in text_new2:
            answer = answer + word + ' '
        print(answer)
        url = f'https://www.google.com/search?q={answer}'
        webbrowser.open(url)
    except:
        voice.speaker('Извините, я не понял запрос! Повторите попытку.')


def start_GPT():
    voice.speaker('включаю для вас искусственный интеллект')
    return True

def stop_GPT():
    voice.speaker('выключаю для вас искусственный интеллект')
    return False


def GPT_answer(text: str):
    print(text)
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        if name == '':
            name = 'Неизвестная личность'

    try:
        with open('API_KEY.txt', 'r', encoding='utf-8') as file:
            api = file.read().strip()
        openai.api_key = api
        chat_completion = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system",
                 "content": f"Твоего пользователя зовут {name}. Ты добрый человек который готов помочь пользователю ответить на любой вопрос. Страрайся давать кратикий ответ на русском языке"},
                {"role": "user",
                 "content": f"Ты - Василий, а твоего пользователя зовут {name}. Запрос: {text}'"}
            ]
        )
        reply = chat_completion.choices.message.content
        voice.speaker(reply)
    except:
        voice.speaker('Ошибка! Проверьте ваш api ключ и соеденение с интернетом!')


def start_GPT_menu():
    voice.speaker('включаю для вас меню с искусственным интеллектом. Пожалуйста подождите')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\my_neyr.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')


def password_server_copy(text: str):
    voice.speaker('сейчас узнаю...')
    text_list = nltk.word_tokenize(text)
    text_list_norm = text_list[6:]
    text_norm = ''
    b = False
    for i in text_list_norm:
        text_norm += i
        text_norm += ' '
    print(text_norm)
    conn = sqlite3.connect('password_server.db')
    cursor = conn.cursor()
    cursor.execute("""SELECT * FROM password""")
    fel = cursor.fetchall()
    for i in fel:
        print(i[0])
        if str(i[0].lower().strip()) in text_norm.lower().strip() or text_norm.lower().strip() == str(i[0]).lower().strip():
            b = True
            pyperclip.copy(str(i[1]))
            voice.speaker(f'Готово! Я смог найти пароль от {text_norm}. Скопировал в буфет обмена. Последняя дата сохранения - {i[2]}')
    if b == False:
        voice.speaker('Извините, я не нашёл пароль. Попробуйте сказать фразу более чётко')

def test_math():
    voice.speaker('запускаю тренажёр')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\math_test.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def music_norm():
    try:
        voice.speaker('Включаю. Приятной игры! Вы пока что слушайте,а я отойду')
        mine()
        pygame.init()
        music_list = ['Agony by Yung Lean instrumental remix.mp3',
                      'analog_mannequin_-_Milk_Cassette_Xmp3_73690486.mp3',
                      'neheart_Reidenshi_-_distorted_memories_75718831.mp3']
        sound_file = random.choice(music_list)
        sound = pygame.mixer.Sound(sound_file)
        sound.set_volume(0.5)
        sound.play()
        while pygame.mixer.get_busy():
            continue
        pygame.quit()
        return False
    except:
        voice.speaker('Ошибка проверьте расположение файла игры или же проверьте наличие песен в папке!')

def music_start(name_music):
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    voice.speaker(random.choice(['ищу музыку', 'начинаю поиск', 'Хорошо, сейчас поищу песню']))
    music_bool = False
    name_music2 = nltk.word_tokenize(name_music)
    file_n = ''
    for i in name_music2[3:]:
        file_n += i
        file_n += ' '
    mus = file_n.strip() + '.mp3'
    print(mus)
    path = f'{path[0][0]}\output'  # НЕ МЕНЯТЬ!!!
    for root, dirs, files in os.walk(path):
        for file in files:
            if file.endswith('.mp3'):
                file_name = file.split('.')
                file_name_new = file_name[0].lower().strip() + '.' + file_name[1]
                print(file_name)
                if mus.lower().strip() == file_name_new:
                    try:
                        voice.speaker('Нашёл! Приятного слушания, ну а я пока что отойду за кружкой чая')
                        pygame.init()
                        sound_file = file
                        sound = pygame.mixer.Sound(sound_file)
                        sound.set_volume(0.5)
                        sound.play()
                        while pygame.mixer.get_busy():
                            continue
                        pygame.quit()
                        music_bool = True
                        return False
                    except:
                        voice.speaker('Извините не нашёл песню проверьте наличие песни в папке output!')
    if music_bool == False:
        voice.speaker('Извините, я не смог найти нужную песню. Попробуйте сказать чётко. Также название должно быть на русском без спец. символов и знаков препинания. Также проверьте наличие песни в папке output!')
        return False

def idea():
    list_idea = ['поиграть', 'почитать', 'погулять', 'заняться спортом', 'сделать вещь своими руками', 'приготовить что-то', 'сходить в музей',
                 'порисовать','сделать генеральную уборку', 'потренироваться в любимом деле', 'написать книгу', 'спеть в какаоке','сходить в боулинг с другом',
                 'погулять с домашним питомцем', 'сходить в гости к другу', 'пообщаться с друзьями', 'сделать работу или дз', 'найти работу']
    idea_n = random.choice(list_idea)
    voice.speaker("Вы можете" + idea_n)

def create_folder(name_n):
    voice.speaker('создаю папку')
    try:
        name1 = nltk.word_tokenize(name_n)
        name2 = ''
        for i in name1[5:]:
            name2 += i
            name2 += ' '
        name = name2.strip()
        print(name)
        desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
        file_name = name
        file_path = os.path.join(desktop_path, file_name)
        os.makedirs(file_path)
        voice.speaker('готово! папка была успешно создана')
    except:
        voice.speaker('Ошибка создания! Пожалуйста проверьте название или скажите мне чётко')

def create_text_file(name_n):
    voice.speaker('создаю текстовый файл')
    try:
        name1 = nltk.word_tokenize(name_n)
        name2 = ''
        for i in name1[6:]:
            name2 += i
            name2 += ' '
        name = name2.strip()
        desktop_path = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop')
        file_name = name + '.txt'

        with open(os.path.join(desktop_path, file_name), "w") as file:
            file.write(" ")
            voice.speaker('готово! текстовый файл успешно создан')
    except:
        voice.speaker('Ошибка создания! Пожалуйста проверьте название или скажите мне чётко')

def clear_basket():
    voice.speaker('начинаю очистку корзины')
    try:
        winshell.recycle_bin().empty(confirm=False, show_progress=False, sound=True)
        voice.speaker('Корзина была успешно очищена!')
    except:
        voice.speaker('Ошибка очистки! Попробуйте ещё раз. Возможно в корзине ничего нет!')

def open_blustacks():
    voice.speaker('открываю')
    try:
        subprocess.Popen("C:\Program Files (x86)\BlueStacks X\BlueStacks X.exe")
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def rand():
    r = random.randint(1,100)
    voice.speaker(f'выпало число: {r}')

def you_video(text):
    voice.speaker('открываю ютуб')
    text_list = nltk.word_tokenize(text)
    text_word = text_list[5:]
    text_norm = ''
    for i in text_word:
        text_norm += i
        text_norm += ' '
    text_norm.strip()
    print(text_norm)
    url = f'https://www.youtube.com/results?search_query={text_norm}'
    webbrowser.open(url)

def text_itog():
    list_answer = ['да', 'нет', 'возможно', 'не желательно', 'скорее да чем нет', 'скорее нет чем да']
    r = random.choice(list_answer)
    voice.speaker(f'Мой ответ: {r}')

def yandex_music():
    voice.speaker('открываю яндекс музыку')
    webbrowser.open('https://music.yandex.ru/home?t')

def python_file_create(text):
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    text_new = nltk.word_tokenize(text)
    text_new_2 = text_new[6:]
    text_norm = ''
    for i in text_new_2:
        text_norm += i
        text_norm += ' '
    text_norm = text_norm.strip()
    file_name = text_norm + '.py'

    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
    with open(f'{path[0][0]}\рабочие файлы\{file_name}', 'w', encoding='utf-8') as file2:
        file2.write(f"#Добро пожаловать в python. Приятной работы {name.capitalize()}!\n")
    voice.speaker('готово! я создал пайтон файл в папке: рабочие файлы. Приятной работы!')


def interesting_facts():
    list_facts = ['В среднем одинокие мужчины на два с половиной сантиметра ниже ростом, чем женатые.',
                  'Виноград взрывается в микроволновой печи.',
                  'У блондинов больше волос.',
                  'Солнечный свет достигает Земли за 8 минут, Юпитера — за 40 минут, а чтобы достигнуть границ Солнечной системы, ему требуется 7 часов.',
                  'Без головы таракан может жить еще 2 недели!',
                  'Крыса может упасть с пятиэтажного здания без каких-либо повреждений.',
                  'Футбольный стадион Уэмбли в Англии создан только для игр, тренировки на нем запрещены.',
                  'Температура крови у рыб Антарктиды может достигать -1,7 градусов Цельсия.',
                  'Аляска — Самый северный, Восточный, И Западный штат во всей Америке. Это также единственная часть континента, входящая в восточное полушарие.',
                  'Синий кит - самое крупное млекопитающее, весящее столько же, сколько 23 слона.',
                  'Глаз страуса больше, чем его мозг.',
                  'В Антарктиде есть водопад с красной, как кровь, водой, что объясняется наличием железа, которое окисляется при контакте с воздухом.',
                  'На территории Антарктиды есть семь христианских церквей.',
                  '99% площади Антарктиды находится подо льдом. Ледник, покрывающий этот континент, называется «ледяным щитом».',
                  'Самый большой динозавр из когда-либо обнаруженных достигал в длину более тридцати метров и весил более восьмидесяти тонн.',
                  'Продолжительность первого выхода в космос Леоновым составила двенадцать минут.',
                  'Игрушка йо-йо появилась в шеснадцатом веке на Филиппинах в качестве оружия.',
                  'Национальный гимн Греции имеет 158 куплетов.',
                  'Из двух миллиардов человек, только одному удается дожить до 116 лет.',
                  'Каждый год Солнце «теряет» 360 миллионов тонн.',
                  'Самая короткая война в истории – между Англией и Занзибаром, продолжалась всего 38 минут.',
                  'Раньше в разработке голосовой помощник Василий имел лицо и тело.']
    random_fact = random.choice(list_facts)
    voice.speaker(f'Интересный факт: {random_fact}')

def disk_D():
    voice.speaker('сейчас узнаю')
    try:
        disk_usage = psutil.disk_usage('D:/')
        total_gb = disk_usage.total / (2 ** 30)
        gb_used1 = disk_usage.used / (2 ** 30)
        gb_used = total_gb - gb_used1
        number = (gb_used1 / total_gb) * 100
        voice.speaker(
            f'Всего на диске дэ используется {round(total_gb)} гб, используется {round(gb_used)}. Всего используется {round(number)} процент!')
    except:
        voice.speaker('Ошибка!')

def disk_C():
    try:
        voice.speaker('сейчас узнаю')
        disk_usage = psutil.disk_usage('C:/')
        total_gb = disk_usage.total / (2 ** 30)
        gb_used1 = disk_usage.used / (2 ** 30)
        number = gb_used1 / total_gb * 100
        voice.speaker(f'Всего на диске цэ используется {round(total_gb)} гб, используется {round(gb_used1)}. Всего используется {round(number)} процент!')
    except:
        voice.speaker('Ошибка!')

def info_user():
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
    conn = sqlite3.connect("user_info.db")
    cursor = conn.cursor()
    cursor.execute("""SELECT * FROM user""")
    all = cursor.fetchall()
    if str(all) == '[]' or name.strip() == '':
        voice.speaker('Вы не вносили о себе информацию! Чтобы я смог что-то вам рассказать внесите о себе информацию в главном меню')
    else:
        age = all[0][0]
        mass = all[0][1]
        height = all[0][2]
        hobby = all[0][3]
        city = all[0][4]
        voice.speaker(f'Краткое инфо про вас: вас зовут {name}. Ваше хобби на данный моменты - это {hobby}. Сейчас вам {age} лет, ваш вес состовляет {mass} килограмм и ваш рост состовляет {height} см. Сейчас вы проживаете в городе {city}')


def password_menage():
    voice.speaker('открываю менеджер паролей')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\password_message.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')



def start_GPT_talk():
    voice.speaker('хорошо, давайте поговорим')
    return True

def GPT_talk(text: str):
    voice.speaker('Сейчас отвечу на ваш вопрос!')
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        if name == '':
            name = 'Неизвестная личность'
    try:
        with open('API_KEY.txt', 'r', encoding='utf-8') as file:
            api = file.read().strip()
        openai.api_key = api
        chat_completion = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system",
                 "content": f"Твоего пользователя зовут {name}, а тебя Василий(Всегда василий, не название модели!). Ты создан только для общения с {name}. Старайся давать ему краткие ответы, иногда можно длинные. Твоя главная цель - это общение с {name}"},
                {"role": "user",
                 "content": f"Ты - Василий, а твоего пользователя зовут {name}. Запрос: {text}'"}
            ]
        )
        reply = chat_completion.choices.message.content
        voice.speaker(reply)
    except:
        voice.speaker('Ошибка! Проверьте ваш api ключ и соеденение с интернетом!')
def GPT_talk_stop():
    voice.speaker('Хорошо, отлично поговорили!')
    return False

def meet():
    meet_l = ['всё отлично', 'всё нормально', 'всё прекрасно', 'всё хорошо', 'вполне неплохо']
    r_meet = random.choice(meet_l)
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        voice.speaker(f'{name} у меня {r_meet}, а у вас?')

def open_kalk():
    voice.speaker('открываю калькулятор')
    try:
        cursor_path.execute("""SELECT * FROM your_path""")
        path = cursor_path.fetchall()
        subprocess.Popen(f'{path[0][0]}\output\calk.exe')
    except:
        voice.speaker('Ошибка! Проверьте расположения файла')

def get_new():
    voice.speaker(random.choice(['Узнаю новость', 'Начинаю поиск', 'Нахожу новость']))
    news_topics = ["Политика", "Экономика", "Технологии", "Спорт", "Здоровье",
                   "Культура", "Наука", "Образование", "Игры", "Путешествия"]

    data = str(datetime.datetime.now())[:10]  # Дата за сегодня
    yesterday = datetime.date.today() - datetime.timedelta(days=1)  # Дата за прошлый день

    url = ('https://newsapi.org/v2/everything?'
           'q=' + random.choice(news_topics) + '&' # Тема новости
           'from=' + str(data) + '&' #За какое число получить новость
           'sortBy=popularity&' #Сортировка новостей
           'apiKey=abffbb74685040f0acb39c1722402ad8') #apikey нужно сгенерировать индивидуально!
    response = requests.get(url) # Связь с API
    try:
        answer = response.json()['articles'][0]['title'] + '. ' + response.json()['articles'][0]['description']
        url = response.json()['articles'][0]['url']
    except IndexError as e:
        url = ('https://new'
               'sapi.org/v2/everything?'
               'q=' + random.choice(news_topics) + '&'
               'from=' + str(yesterday) + '&'
               'sortBy=popularity&'
               'apiKey=abffbb74685040f0acb39c1722402ad8')
        response = requests.get(url)  # Связь с API
        answer = response.json()['articles'][0]['title'] + '. ' + response.json()['articles'][0]['description']
        url = response.json()['articles'][0]['url']

    pyperclip.copy(url) #Копирование ссылки на новость
    voice.speaker(answer + '. Для полного ознакомления с новостью можете перейти по скопированной ссылке.')


def weather_in_your_city():
    voice.speaker(random.choice(['узнаю погоду в вашем городе', 'начинаю поиск', 'ищу вашу погоду']))
    conn = sqlite3.connect("user_info.db")
    cursor = conn.cursor()
    cursor.execute("""SELECT * FROM user""")
    all = cursor.fetchall()
    if str(all) == '[]':
        city = 'неизвестный город'
    else:
        city = all[0][4]
    print(city)
    try:
        time = datetime.datetime.now()
        if time.minute < 10:
            minute = '0' + str(time.minute)
        else:
            minute = str(time.minute)
        if time.hour < 10:
            hour = '0' + str(time.hour)
        else:
            hour = str(time.hour)

        time_new = hour + ':' + minute
        conn2 = sqlite3.connect('wether_data.db')
        cur = conn2.cursor()
        cur.execute("""SELECT * FROM weather""")
        all_weather = cur.fetchall()
        print(len(all_weather))
        if len(all_weather) >= 10:
            cur.execute("""DELETE FROM weather""")
            conn2.commit()

        current_datetime = datetime.datetime.now()
        moun = str(current_datetime).split()[0]


        url = 'https://api.openweathermap.org/data/2.5/weather?q=' + city + '&units=metric&lang=ru&appid=79d1ca96933b0328e1c7e3e7a26cb347'
        weather_data = requests.get(url).json()
        temp = round(weather_data['main']['temp'])
        weather_status = weather_data['weather'][0]['description']

        if int(temp) <= 10:
            voice.speaker(
                f'Сейчас в городе {city} температура состовляет {round(temp)} градус цельсия. Лучше оденьтесь потеплее. Тип погоды: {weather_status}')
        if int(temp) <= 20 and int(temp) >= 11:
            voice.speaker(
                f'Сейчас в городе {city} температура состовляет {round(temp)} градус цельсия. Неплохая погода для прогулки. Тип погоды: {weather_status}')
        elif int(temp) > 20:
            voice.speaker(
                f'Сейчас в городе {city} температура состовляет {round(temp)} градус цельсия. Сегодня будет очень жарко. Тип погоды: {weather_status}')
        cur.execute("""INSERT INTO weather VALUES (?,?,?,?,?)""",
                    (str(city), str(moun), str(time_new), temp, str(weather_status)))
        conn2.commit()
    except:
        voice.speaker('Извините я не смог найти погоду. Пожалуйста проверьте вводили ли вы свои данные или же проверьте корректно ли вы ввели данные! Также проверте соединение с интернетом')


def set_language():
    cursor_path.execute("""SELECT * FROM your_path""")
    path = cursor_path.fetchall()
    subprocess.Popen(f'{path[0][0]}\output\othen_language.exe')


def presentation(text):
    text2 = nltk.word_tokenize(text)
    text_norn = text2[5:]
    t = ''
    for i in text_norn:
        t += i
        t += ' '
    t2 = t.strip()
    voice.speaker(random.choice(['Создаю', 'Начинаю создание', 'Генерирую пустую презентацию']))
    with open(f'{path[0][0]}\my_name', 'r', encoding='utf-8') as file:
        name = file.read()
        root = Presentation()

        first_slide = root.slide_layouts[0]
        slide = root.slides.add_slide(first_slide)
        slide.shapes.title.text = 'Добро пожаловать!'
        slide.placeholders[1].text = f'приятной работы\n{name}'

        root.save(f'{path[0][0]}\презентации\{t2}.pptx')
        voice.speaker('Готово! Я создал презентацию в папке презентации!')

def open_generate():
    subprocess.Popen(f'{path[0][0]}\output\generate.exe')
    voice.speaker('Открываю генератор презентаций')

def stop_listen():
    voice.speaker(random.choice(['Прекращаю прослушивание', 'завершаю работу на время', 'Хорошо, жду вас']))
    return False

def start_listen():
    voice.speaker(random.choice(['Возобнавляю прослушивание', 'Я вас снова слушаю', 'Ещё раз здравствуйте', 'Я вернулся']))
    return True

def open_reading():
    voice.speaker('Открываю мини читалку ткст файлов')
    subprocess.Popen(f'{path[0][0]}/output/reading.exe')


def rock():
    voice.speaker('Открываю игру в камень ножницы бумагу')
    subprocess.Popen(f'{path[0][0]}/output/rock_paper_scissors.exe')

def gallows():
    voice.speaker('Открываю виселицу')
    subprocess.Popen(f'{path[0][0]}/output/gallows.exe')

def last_zam():
    with open(f'{path[0][0]}/заметки и шифры/записки.txt', 'r', encoding='utf-8') as file:
        t = file.readlines()
        answer = t[-2]
        if answer[-1] == '\n':
            answer = t[-3]

    voice.speaker(f'Последняя {answer}')

def play_music():
    voice.speaker('Открываю проигрыватель')
    subprocess.Popen(f"{path[0][0]}/output/music_player.exe")


def set_mus(volume_num: int) -> None:
    if not volume_num: voice.speaker('Выключаю звук')
    else: voice.speaker('Включаю звук')
    devices = AudioUtilities.GetSpeakers()
    interface = devices.Activate(
        IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
    volume = interface.QueryInterface(IAudioEndpointVolume)
    current_volume = volume.GetMasterVolumeLevelScalar()
    print(f"Текущая громкость: {current_volume}")
    volume.SetMasterVolumeLevelScalar(volume_num, None)
    print(f"Громкость установлена на: {volume_num}")

def close_mus():
    new_volume = 0
    set_mus(new_volume)


def open_mus():
    new_volume = 1
    set_mus(new_volume)