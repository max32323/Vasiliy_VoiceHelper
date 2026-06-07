"""Импорт нужных библиотек"""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random as rnd
import image_get
import os
import get_text
import voice
import shutil
import sqlite3
import time

connect_path = sqlite3.connect('user_path.db')
cursor_path = connect_path.cursor()
cursor_path.execute("""SELECT * FROM your_path""")
path = cursor_path.fetchall()


"""Генерация самой презентации"""
def generate_presintation(name_presentation,aftor,plan, file2, save):
    try:
        save_in_table = save

        photo_b = True
        voice.speaker('Генерация началась. Пожалуйста подождите!')
        """Создание"""
        prs = Presentation()

        """Дизайн(цвет слайда)"""
        while True:
            green = rnd.randint(1, 255)
            red = rnd.randint(1, 255)
            blue = rnd.randint(1, 255)
            if (green != red != blue):
                break
        f = rnd.randint(1, 2)
        f2 = rnd.randint(1, 2)
        new_color = RGBColor(red, green, blue)

        slide_layout = prs.slide_layouts[5]

        list_font = ['Microsoft JhengHei UI Light', 'Times New Roman', 'Franklin Gothic Medium', 'Bahnschrift Light',
                     'Segoe UI Semibold', 'Yu Gothic UI Light', 'Century',
                     'Century Gothic', 'Trebuchet MS', 'Cascadia Code SemiLight', 'Impact', 'Ink Free', 'PMingLiU-ExtB', 'Lucida Console']

        list_font2 = ['Century', 'Arial Narrow', 'Arial', 'Cambria Math', 'Microsoft New Tai Lue', 'Georgia', 'Candara Light', 'Gabriola', 'Segoe Print', 'MV Boli', 'Comic Sans MS']

        font = rnd.choice(list_font)

        font2 = rnd.choice(list_font2)

        list_makets = [('maket1.PNG', 'maket1_2.PNG'), ('maket2.PNG', 'maket2_2.PNG'), ('maket3_2.PNG', 'maket3_2.PNG'),
                       ('maket4.PNG', 'maket4_2.PNG'), ('maket5.PNG', 'maket5_2.PNG'),
                       ('maket6.PNG', 'maket6_2.PNG'), ('maket7.jpg', 'maket7.jpg'), ('maket8.jpg', 'maket8.jpg'),
                       ('maket9.jpg', 'maket9.jpg'), ('maket10.jpg', 'maket10.jpg'),
                       ('maket11.PNG', 'maket11.PNG'), ('maket12_1.PNG', 'maket12_2.PNG'), ('maket13_1.PNG', 'maket13_2.PNG')]

        bool_makets = rnd.randint(0, 4)

        maket = rnd.choice(list_makets)
        #maket = ('maket7.jpg', 'maket7_2.jpg')
        print(bool_makets)
        print(maket)

        list_s = ['◆', '∎', '', '✮', '▶']
        sym = rnd.choice(list_s)

        number_check = rnd.choice([True, False])

        """Создание слайдов"""
        for i in range(len(plan) + 1):
            slide = prs.slides.add_slide(slide_layout)

        count = 0
        print(plan, '\n')
        rand = rnd.randint(1, 4)
        number_turn = 1
        print('ЭТАП 1 ПРОЙДЕН')

        for slide in prs.slides:
            count += 1
            # print(plan[count-1])
            if (count == 1):
                """Слайд начало"""

                if (bool_makets == 0):
                    background = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, left=0, top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    background.fill.solid()
                    background.fill.fore_color.rgb = new_color
                else:
                    slide.shapes.add_picture(maket[0], Inches(0), Inches(0), width=Inches(10))

                left = Inches(4.1)
                top = Inches(0.5)
                width = Inches(5)
                height = Inches(1)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame

                p = text_frame.add_paragraph()
                p.text = "Тема: "
                p.font.size = Pt(30)
                p.font.name = 'Arial Black'

                # Заголовок
                left = Inches(0.9)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(1)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame

                g = 0

                print('ЭТАП 2 ПРОЙДЕН')

                if (len(name_presentation) <= 40):
                    p = text_frame.add_paragraph()
                    p.text = '● ' + name_presentation
                    p.font.size = Pt(30)
                else:
                    """Слайды по плану"""
                    chunks1 = []
                    for i in range(0, len(name_presentation), 40):
                        chunk = name_presentation[i:i + 40]
                        chunks1.append(chunk.strip())
                    for i in chunks1:
                        g += 1
                        p = text_frame.add_paragraph()
                        p.text = i
                        p.font.size = Pt(28)
                        p.space_after = Inches(0.01)
                        if (g >= 2):
                            photo_b = False

                    p = text_frame.add_paragraph()
                    p.text = name_presentation[40:].strip()
                    p.font.size = Pt(24)

                left = Inches(0.5)
                top = Inches(6.3)
                width = Inches(5)
                height = Inches(1)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame

                p = text_frame.add_paragraph()
                p.text = f"Автор: {aftor}"
                p.font.size = Pt(22)
                font_aftor = ['Arial Black', 'Arial Unicode MS', 'Yu Gothic UI Semilight',
                              'Bahnschrift SemiBold Condensed', 'Arial Narrow', 'HoloLens MDL2 Assets']
                p.font.name = rnd.choice(font_aftor)
                black = rnd.randint(1, 4)
                if (black == 1):
                    p.font.bold = True
                elif (black == 2):
                    p.font.italic = True
                elif (black == 3):
                    p.font.bold = True
                    p.font.italic = True
                print(maket[0])
                if maket[0] == 'maket2.PNG' or maket[0] == 'maket9.jpg' or maket[0] == 'maket10.jpg' or maket[
                    0] == 'maket7.jpg':
                    pass

                else:
                    if (bool_makets == 0):
                        if (photo_b):
                            img_path = image_get.get_photo(f'Тема презентации: {name_presentation}', 1)
                            img = Image.open(img_path)
                            img_new = img.resize((540, 420))
                            img_new.save(img_path)
                            slide.shapes.add_picture(img_path, Inches(3.2), Inches(3.4), width=Inches(3.5))
                            os.remove(img_path)
                    else:
                        if (photo_b):
                            img_path = image_get.get_photo(f'Тема презентации: {name_presentation}', 1)
                            img = Image.open(img_path)
                            img_new = img.resize((540, 420))
                            img_new.save(img_path)
                            slide.shapes.add_picture(img_path, Inches(3.2), Inches(3.4), width=Inches(3.5))
                            os.remove(img_path)

                prs.save(file2)
                print('ТИТУЛЬНИК СДЕЛАН')
            else:
                print('Тайм стоп')
                time.sleep(8)
                if (bool_makets == 0):
                    background = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, left=0, top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    background.fill.solid()
                    background.fill.fore_color.rgb = new_color
                else:
                    slide.shapes.add_picture(maket[1], Inches(0), Inches(0), width=Inches(10))

                title = slide.shapes.title
                title.text = " "

                left = Inches(0.5)
                top = Inches(0.25)
                width = Inches(5)
                height = Inches(1)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame
                if len(plan[count - 2]) <= 25:

                    b = 0

                    p = text_frame.add_paragraph()
                    p.text = f"{sym} {plan[count - 2]}"
                    p.font.size = Pt(24)
                    p.font.name = font2
                    if (f == 1):
                        p.font.bold = True
                    if (f2 == 1):
                        p.font.italic = True
                else:
                    chunks2 = []
                    t = plan[count - 2]
                    for i in range(0, len(plan[count - 2]), 40):
                        chunk2 = t[i:i + 40]
                        chunks2.append(chunk2.strip())

                    p = text_frame.add_paragraph()
                    p.text = f"{sym} {chunks2[0]}"
                    if (f == 1):
                        p.font.bold = True
                    if (f2 == 1):
                        p.font.italic = True
                    p.font.name = font2
                    p.font.size = Pt(24)
                    p.space_after = Inches(0.01)

                    for i in chunks2[1:]:
                        p = text_frame.add_paragraph()
                        p.text = i
                        if(f==1):
                            p.font.bold = True
                        if (f2 == 1):
                            p.font.italic = True
                        p.font.name = font2
                        p.font.size = Pt(24)
                        p.space_after = Inches(0.01)

                try:
                    img_path = image_get.get_photo(f'{name_presentation}. {plan[count - 2]}')
                    print(f'ПУТЬ: {img_path}')
                    img = Image.open(img_path)
                    img_new = img.resize((540, 420))
                    img_new.save(img_path)
                except:
                    for i in range(image_get.count_photo, 0, -1):
                        try:
                            img = Image.open(img_path.replace(img_path[5], f'{i}'))
                            img_new = img.resize((540, 420))
                            img_new.save(img_path)
                            break
                        except:
                            pass
                if(rand==3):
                    r = rnd.randint(1, 2)
                    if r == 1:
                        slide.shapes.add_picture(img_path, Inches(0.6), Inches(2.4), width=Inches(4.5))

                        left = Inches(5.2)
                        top = Inches(1.2)
                        width = Inches(4)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame
                        p = text_frame.add_paragraph()
                        p.font.size = Pt(16)

                        text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                        print(text_answer)
                        chunks = []
                        for i in range(0, len(text_answer), 30):
                            chunk = text_answer[i:i + 30]
                            chunks.append(chunk.strip())
                        for i in chunks:
                            p = text_frame.add_paragraph()
                            p.text = i
                            p.font.name = font
                            p.font.size = Pt(18)
                            p.space_after = Inches(0.01)
                        if (number_check):
                            left = Inches(0.5)
                            top = Inches(6.3)
                            width = Inches(5)
                            height = Inches(1)

                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = txBox.text_frame

                            p = text_frame.add_paragraph()
                            p.text = f"{str(count - 1)}"
                            p.font.size = Pt(22)
                            p.font.name = font

                    else:
                        slide.shapes.add_picture(img_path, Inches(4.85), Inches(2), width=Inches(4.5))

                        left = Inches(0.4)
                        top = Inches(1.2)
                        width = Inches(4)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame
                        p = text_frame.add_paragraph()
                        p.font.size = Pt(16)

                        text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                        print(text_answer)
                        chunks = []
                        for i in range(0, len(text_answer), 30):
                            chunk = text_answer[i:i + 30]
                            chunks.append(chunk.strip())
                        for i in chunks:
                            p = text_frame.add_paragraph()
                            p.text = i
                            p.font.name = font
                            p.font.size = Pt(16)
                            p.space_after = Inches(0.01)
                        if (number_check):
                            left = Inches(8.8)
                            top = Inches(6.3)
                            width = Inches(5)
                            height = Inches(1)

                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = txBox.text_frame

                            p = text_frame.add_paragraph()
                            p.text = f"{str(count - 1)}"
                            p.font.size = Pt(22)
                            p.font.name = font
                elif(rand==1):
                    slide.shapes.add_picture(img_path, Inches(0.6), Inches(2), width=Inches(4.5))

                    left = Inches(5.2)
                    top = Inches(1.2)
                    width = Inches(4)
                    height = Inches(1)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = txBox.text_frame
                    p = text_frame.add_paragraph()
                    p.font.size = Pt(16)

                    text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                    print(text_answer)
                    chunks = []
                    for i in range(0, len(text_answer), 30):
                        chunk = text_answer[i:i + 30]
                        chunks.append(chunk.strip())
                    for i in chunks:
                        p = text_frame.add_paragraph()
                        p.text = i
                        p.font.name = font
                        p.font.size = Pt(16)
                        p.space_after = Inches(0.01)

                    if(number_check):
                        left = Inches(0.5)
                        top = Inches(6.3)
                        width = Inches(5)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame

                        p = text_frame.add_paragraph()
                        p.text = f"{str(count - 1)}"
                        p.font.size = Pt(22)
                        p.font.name = font
                elif(rand==2):
                    slide.shapes.add_picture(img_path, Inches(4.85), Inches(2), width=Inches(4.5))

                    left = Inches(0.4)
                    top = Inches(1.2)
                    width = Inches(4)
                    height = Inches(1)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = txBox.text_frame
                    p = text_frame.add_paragraph()
                    p.font.size = Pt(18)

                    text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                    print(text_answer)
                    chunks = []
                    for i in range(0, len(text_answer), 30):
                        chunk = text_answer[i:i + 30]
                        chunks.append(chunk.strip())
                    for i in chunks:
                        p = text_frame.add_paragraph()
                        p.text = i
                        p.font.name = font
                        p.font.size = Pt(16)
                        p.space_after = Inches(0.01)

                    if(number_check):
                        left = Inches(8.8)
                        top = Inches(6.3)
                        width = Inches(5)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame

                        p = text_frame.add_paragraph()
                        p.text = f"{str(count - 1)}"
                        p.font.size = Pt(22)
                        p.font.name = font
                elif(rand==4):
                    if number_turn%2 == 0:
                        slide.shapes.add_picture(img_path, Inches(0.6), Inches(2), width=Inches(4.5))

                        left = Inches(5.2)
                        top = Inches(1.2)
                        width = Inches(4)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame
                        p = text_frame.add_paragraph()
                        p.font.size = Pt(18)

                        text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                        print(text_answer)
                        chunks = []
                        for i in range(0, len(text_answer), 30):
                            chunk = text_answer[i:i + 30]
                            chunks.append(chunk.strip())
                        for i in chunks:
                            p = text_frame.add_paragraph()
                            p.text = i
                            p.font.name = font
                            p.font.size = Pt(16)
                            p.space_after = Inches(0.01)

                        if (number_check):
                            left = Inches(0.5)
                            top = Inches(6.3)
                            width = Inches(5)
                            height = Inches(1)

                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = txBox.text_frame

                            p = text_frame.add_paragraph()
                            p.text = f"{str(count - 1)}"
                            p.font.size = Pt(22)
                            p.font.name = font


                    else:
                        slide.shapes.add_picture(img_path, Inches(4.85), Inches(2), width=Inches(4.5))

                        left = Inches(0.4)
                        top = Inches(1.2)
                        width = Inches(4)
                        height = Inches(1)

                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = txBox.text_frame
                        p = text_frame.add_paragraph()
                        p.font.size = Pt(18)

                        text_answer = get_text.GPT_answer(f'ответ по пункту: {plan[count - 2]}. Тема: {name_presentation}. дай краткий ответ до 300 символов.')
                        print(text_answer)
                        chunks = []
                        for i in range(0, len(text_answer), 30):
                            chunk = text_answer[i:i + 30]
                            chunks.append(chunk.strip())
                        for i in chunks:
                            p = text_frame.add_paragraph()
                            p.text = i
                            p.font.name = font
                            p.font.size = Pt(16)
                            p.space_after = Inches(0.01)

                        if (number_check):
                            left = Inches(8.8)
                            top = Inches(6.3)
                            width = Inches(5)
                            height = Inches(1)

                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = txBox.text_frame

                            p = text_frame.add_paragraph()
                            p.text = f"{str(count - 1)}"
                            p.font.size = Pt(22)
                            p.font.name = font
                    number_turn+=1

                prs.save(file2)
                for i in range(1, image_get.count_photo+1):
                    img_path = img_path.replace(img_path[5], str(i))
                    for type in ['jpg', 'png', 'jpeg', 'svg']:
                        if not os.path.exists(img_path):
                            img_path = img_path.replace(img_path[7:], type)
                        else:
                            os.remove(img_path)

        if rnd.randint(1,2) == 1:
            slide = prs.slides.add_slide(slide_layout)
            list_end = ['end1.jpg', 'end2.jpg', 'end3.jpg', 'end4.jpg', 'end5.jpg', 'end6.jpg', 'end7.png']
            end = rnd.choice(list_end)
            slide.shapes.add_picture(end, Inches(0), Inches(0), width=Inches(10))
            prs.save(file2)
            print('Готово')
        move(file2, save_in_table)
    except Exception as e:
        print(e)
        voice.speaker('Ошибка генерации. Повторите попытку позже. Также все окна повер поинт должны быть закрыты!')


def move(name_file, save_in_table):
    source_file = f'{path[0][0]}\output\{name_file}'

    #Сохранение в спец. папку
    if save_in_table == False:
        desktop_path = f'{path[0][0]}\презентации'
        target_file = os.path.join(desktop_path, os.path.basename(source_file))
        text = 'Презентация была сгенерирована! Её можно найти в папке презентации!'

    #Сохранение на рабочий стол
    else:
        desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
        target_file = os.path.join(desktop_path, f'{name_file}')
        text = 'Презентация была сгенерирована! Её можно найти на рабочем столе!'

    try:
        shutil.move(source_file, target_file)
        voice.speaker(text)
    except PermissionError:
        voice.speaker('Ошибка перемещения! Пожалуйста проверьте разрешения перемещения в указанную область!')

#generate_presintation('Роботы', 'Максим', ['Введение','Что такое робот?','Применение','Вывод'], 'ex1.pptx', False)